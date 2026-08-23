import os
from pptx import Presentation
from pptx.util import Inches, Pt

def inject_text_into_slide(slide, body_text):
    # Find the main text shape (skip shapes that are purely decorative)
    main_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text.strip():
            # Usually the main shape is the one with the most text
            main_shape = shape
            break
            
    if not main_shape:
        return
        
    tf = main_shape.text_frame
    
    # Keep paragraph 0 (title) intact. Delete all other paragraphs.
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)
    
    # Add a blank line for spacing
    p_blank = tf.add_paragraph()
    p_blank.text = ""
    
    # Add our new dense body text
    p_body = tf.add_paragraph()
    p_body.text = body_text
    p_body.font.size = Pt(18)  # Professional size that won't overlap

def add_image_to_slide(slide, image_path, left_in, top_in, width_in):
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in), width=Inches(width_in))
        else:
            print(f"Warning: Image {image_path} not found.")
    except Exception as e:
        print(f"Error adding image {image_path}: {e}")

def main():
    prs = Presentation('[EXT] UniHack-Protoype Template  (1).pptx')
    
    # 1. DELETE FIRST SLIDE (Guidelines)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])
    # Now old Slide 1 (Team Details) is index 0
    
    # Team Details Slide
    tf = prs.slides[0].shapes[1].text_frame
    tf.paragraphs[2].text = "Team name: KLassic"
    tf.paragraphs[3].text = "Team leader name: Abhijeet Kangane"
    
    # Slide 1: Brief about your solution
    inject_text_into_slide(prs.slides[1], "OmniSpec is a comprehensive AI-powered intelligence pipeline that bulk processes unstructured supplier datasheets (PDFs) and URLs into perfectly validated, commerce-ready catalogs. It leverages advanced Large Language Models (LLMs) and deterministic fallback engines to generate structured Golden Records with extremely high accuracy, bypassing the need for manual data entry.")
    
    # Slide 2: Questions (Enrichment, Validation, Scalability)
    # This slide has a multiline title. We will just completely replace the text safely.
    tf = prs.slides[2].shapes[1].text_frame
    while len(tf.paragraphs) > 0:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)
    
    p = tf.add_paragraph()
    p.text = "1. Enrichment: Ingests minimal inputs and fetches technical sheets, using a Hybrid AI extraction pipeline (Qwen-27B LLM + Deterministic Fallback).\n\n2. Validation: Applies real-time Confidence Scoring. Anything below 95% is flagged for Human-in-the-Loop (HITL) review to ensure 100% integrity.\n\n3. Scalability: Built on FastAPI and Next.js microservices for asynchronous bulk processing. Uses PyMuPDF/BS4 to dynamically normalize any unstructured document format."
    p.font.size = Pt(16)

    # Slide 3: Opportunities
    inject_text_into_slide(prs.slides[3], "Unlike rigid parsers, OmniSpec dynamically adapts to any manufacturer's layout automatically. \n\nUnlike pure API-dependent LLMs which crash frequently, our hybrid engine guarantees 100% reliable local fallback during network outages. \n\nUSP: Zero-API-Dependency Fallback Mode for 100% uptime, Ultra-Fast asynchronous processing, and an automated UNSPSC taxonomy prediction engine natively integrated into a premium sci-fi dashboard.")
    
    # Slide 4: Features
    inject_text_into_slide(prs.slides[4], "• Multi-modal Ingestion (Drag-and-Drop PDFs, Direct URLs)\n• Hybrid AI Extraction Pipeline (Groq LLM + Local Deterministic Fallback)\n• Automated UNSPSC Taxonomy Categorization\n• Real-time Confidence Scoring and Data Triangulation\n• Human-in-the-Loop (HITL) Enrichment Dashboard\n• Fully Responsive Sci-Fi UI for industrial users")
    
    # Slide 8: Technologies Used
    inject_text_into_slide(prs.slides[8], "Frontend: Next.js 16.3, React 19, Tailwind CSS, Framer Motion\nBackend: FastAPI, Uvicorn, Python 3.10+\nAI/Parsing: Groq API (Qwen-27B), PyMuPDF, BeautifulSoup4\nDeployment: Vercel (Frontend), Render (Backend)\nTooling: Git, Eslint")
    
    # Slide 9: Cost
    inject_text_into_slide(prs.slides[9], "Extremely Low / Near Zero.\n\nBy utilizing high-efficiency open-source models (Llama3/Qwen) via ultra-fast inference APIs (Groq), the operational processing cost per 10,000 SKUs is a fraction of traditional manual data entry or commercial LLM (GPT-4) alternatives.")
    
    # Slide 11: Future Dev
    inject_text_into_slide(prs.slides[11], "Our immediate future roadmap includes:\n\n1. Direct eCommerce platform integrations (e.g., Shopify, Magento ERP plugins).\n2. An automated web-crawler that actively monitors supplier websites for real-time datasheet updates and pushes delta changes directly to the Golden Records.")
    
    # Slide 12: Links
    tf = prs.slides[12].shapes[1].text_frame
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._element.getparent().remove(p._element)
    
    p = tf.add_paragraph()
    p.text = "\nGitHub: https://github.com/Abhi666-max/OmniSpec\nDemo Video: https://youtu.be/775Z2EbCy8E\nPrototype Link: https://omni-spec-one.vercel.app/"
    p.font.size = Pt(18)

    # INSERT GENERATED IMAGES
    # Perfectly centered, top=1.8 (below title), width=6.5 (height=3.65) so bottom is 5.45 (safe above blue line)
    # Left = (10 - 6.5) / 2 = 1.75
    img_left, img_top, img_width = 1.75, 1.8, 6.5
    
    # Slide 5 - Process Flow
    inject_text_into_slide(prs.slides[5], "") # clear placeholder instructions
    add_image_to_slide(prs.slides[5], r"C:\Users\abhij\.gemini\antigravity-ide\brain\de377572-b551-49fe-b84b-c465213fec1f\process_flow_diagram_1787479221560.jpg", left_in=img_left, top_in=img_top, width_in=img_width)
    
    # Slide 6 - Wireframes (Landing Page)
    inject_text_into_slide(prs.slides[6], "") 
    add_image_to_slide(prs.slides[6], r"C:\Users\abhij\.gemini\antigravity-ide\brain\de377572-b551-49fe-b84b-c465213fec1f\landing_page_1787479256961.png", left_in=img_left, top_in=img_top, width_in=img_width)
    
    # Slide 7 - Architecture
    inject_text_into_slide(prs.slides[7], "") 
    add_image_to_slide(prs.slides[7], r"C:\Users\abhij\.gemini\antigravity-ide\brain\de377572-b551-49fe-b84b-c465213fec1f\architecture_diagram_1787479236567.jpg", left_in=img_left, top_in=img_top, width_in=img_width)
    
    # Slide 10 - Snapshots MVP (Dashboard)
    inject_text_into_slide(prs.slides[10], "") 
    add_image_to_slide(prs.slides[10], r"C:\Users\abhij\.gemini\antigravity-ide\brain\de377572-b551-49fe-b84b-c465213fec1f\dashboard_1787479335882.png", left_in=img_left, top_in=img_top, width_in=img_width)

    output_pptx = os.path.abspath("Final_OmniSpec_Presentation.pptx")
    prs.save(output_pptx)
    print(f"Saved PPTX to {output_pptx}")
    
    # Export to PDF
    pdf_path = os.path.abspath("Final_OmniSpec_Presentation.pdf")
    print(f"Exporting to PDF: {pdf_path}")
    try:
        import comtypes.client
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(output_pptx)
        presentation.SaveAs(pdf_path, 32)
        presentation.Close()
        powerpoint.Quit()
        print("PDF Export complete! Saved to:", pdf_path)
    except Exception as e:
        print("Failed to export to PDF automatically (Ensure Microsoft PowerPoint is installed). Error:", str(e))

if __name__ == "__main__":
    main()
